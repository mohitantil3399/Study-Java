import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

# --- BRAND CONFIG (from HTML) ---
COLOR_BG = RGBColor(253, 251, 247)      # #FDFBF7
COLOR_INDIGO = RGBColor(30, 27, 75)     # #1E1B4B
COLOR_GOLD = RGBColor(217, 119, 6)      # #D97706

# --- HELPERS ---

def set_slide_background(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_title(slide, text):
    title = slide.shapes.title
    title.text = text
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.name = "Times New Roman"  # Proxy for Playfair Display
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_INDIGO
        paragraph.alignment = PP_ALIGN.LEFT

def add_content_text(slide, text_content, font_size=18, is_bullet=True):
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.word_wrap = True

    first = True
    for line in text_content:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.font.name = "Arial"  # Proxy for Work Sans
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_INDIGO
        if line.startswith("##"):
            p.text = line.replace("##", "").strip()
            p.font.bold = True
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = COLOR_GOLD
            p.level = 0
        else:
            p.level = 1 if is_bullet else 0

def create_color_swatch(slide, left, top, color, label):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.5), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = COLOR_INDIGO

    tb = slide.shapes.add_textbox(Inches(left), Inches(top + 1.6), Inches(1.5), Inches(0.5))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Arial"
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_INDIGO
    p.alignment = PP_ALIGN.CENTER

def add_flow_step(slide, text, left, top, width=2.6, height=1.0, gold=True):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_GOLD if gold else COLOR_INDIGO
    box.line.color.rgb = COLOR_INDIGO

    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Times New Roman"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_BG
    p.alignment = PP_ALIGN.CENTER
    return box

def add_arrow(slide, left, top, width=0.8, height=0.4):
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_INDIGO
    arrow.line.color.rgb = COLOR_INDIGO
    return arrow

def add_caption(slide, text, left, top, width=4.5, height=0.6):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_INDIGO
    p.alignment = PP_ALIGN.LEFT
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    return tb

def add_image(slide, img_path, left, top, width=None, height=None):
    if not img_path or not os.path.exists(img_path):
        return None
    pic = slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                   width=Inches(width) if width else None,
                                   height=Inches(height) if height else None)
    return pic

def add_sticker_label(slide, text, left, top, bg_color=COLOR_GOLD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CLOUD, Inches(left), Inches(top), Inches(2.0), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = COLOR_INDIGO
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_BG
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_bar_chart_mock(slide, left, top, bars, labels, max_height=2.0, width_per_bar=0.4):
    """Simple bar chart mock using rectangles; bars is list of values 0..1"""
    for i, val in enumerate(bars):
        bar_height = max(0.1, val) * max_height
        x = left + i * (width_per_bar + 0.2)
        y = top + (max_height - bar_height)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width_per_bar), Inches(bar_height)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLOR_GOLD if i % 2 == 0 else COLOR_INDIGO
        rect.line.color.rgb = COLOR_INDIGO
        lb = slide.shapes.add_textbox(Inches(x - 0.05), Inches(top + max_height + 0.05), Inches(width_per_bar + 0.2), Inches(0.4))
        tf = lb.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = labels[i]
        p.font.name = "Arial"; p.font.size = Pt(10); p.font.color.rgb = COLOR_INDIGO
        p.alignment = PP_ALIGN.CENTER

# --- MAIN GENERATOR ---

def generate_presentation(
    output_path="Emami_Sattva_Rebrand_Round1_Expanded.pptx",
    team_name="[Your Team Name]",
    team_members="[Member 1], [Member 2], [Member 3]",
    hero_img=None,
    collection_img=None,
    traceability_img=None
):
    prs = Presentation()

    # SLIDE 1: COVER
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "REBRANDING WARS: ROUND 1 ✨"
    subtitle.text = (
        f"Team Name: {team_name}\n"
        f"Team Members: {team_members}\n"
        "Category: Website Home Page Redesign\n"
        "Brand: Emami (Rebranded as Emami Sattva) 🌿"
    )
    title.text_frame.paragraphs[0].font.name = "Times New Roman"
    title.text_frame.paragraphs[0].font.size = Pt(42)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_INDIGO
    title.text_frame.paragraphs[0].font.bold = True
    for p in subtitle.text_frame.paragraphs:
        p.font.name = "Arial"
        p.font.color.rgb = COLOR_GOLD
        p.font.bold = True
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.LEFT
    add_sticker_label(slide, "Premium • Heritage • Ayurveda", 8.0, 0.8)

    # SLIDE 2: DIAGNOSIS (350+ words)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide)
    add_title(slide, "🏺 The Problem: A Legacy Brand in a Changing World")

    slide2 = [
        "## The Perception Gap ⚠️",
        "Emami is synonymous with household wellness across India, with iconic sub-brands like Zandu and Navratna anchoring its broad reach. Yet, that very ubiquity has created a paradox: in the emerging premium D2C wellness space, ubiquity can be misread as commodity. The modern consumer—especially Gen-Z and Millennials—seeks brands that embody a lifestyle, not just a product. They gravitate toward experiences that feel intentional, ritualistic, and transparent. On our site, the hero promise “Taste the Wisdom of India” hints at this ritualistic aspiration, but the old perception remains: utility-first rather than experience-first.",
        "## The Market Void 📉",
        "While the ghee category is crowded, it is predominantly transactional and commoditized. Competing brands often emphasize price, fat percentages, or generic purity claims without elevating the act of consumption into a daily ritual. Our HTML framework positions Emami Sattva as a bridge: a heritage infusion series built on the Vedic Bilona method and infused botanicals, transforming ghee from a kitchen staple into a wellness instrument. This addresses a notable void: Ritualistic Wellness—products that harmonize culinary tradition with therapeutic Ayurveda, rooted in transparent storytelling. The visual direction (floating turmeric and garlic ingredients) explicitly counters commodity staging by making raw purity the central visual metaphor.",
        "## UX Friction 🧱",
        "Legacy e-commerce patterns tend to hard-pivot to SKU grids, discounts, and multi-click checkout paths, often neglecting narrative arcs and emotional onboarding. Users encounter cluttered navbars, dense product tiles, and inflexible filters, which depress engagement and reduce trust. Our site counters this with a sticky navbar that changes transparency on scroll, subtle hover micro-interactions, and content sections that unfold a journey. Still, the diagnosis is clear: Emami’s digital presence must prioritize calm, whitespace, and story-driven immersion. The cream background (#FDFBF7) provides room to breathe; Indigo (#1E1B4B) builds trust; Gold (#D97706) signals warmth and prosperity without shouting.",
        "## The Strategic Verdict 🧘",
        "To win Round 1, the brand must pivot from selling “oil” to selling “Prana”—from commodity to a life-affirming ritual. That requires a UX philosophy that values pace, ceremony, and clarity: Attract through evocative hero imagery; Educate with heritage and process; Create Desire through infusions; Cement Trust via traceability; and drive Action with obvious, low-friction CTAs. This reframe is not superficial; it transforms Emami’s mass-market strength into aspirational credibility by elevating daily nourishment into mindful practice. In short: reposition the brand as a premium ritual partner, not merely a pantry item."
    ]
    add_content_text(slide, slide2, font_size=18, is_bullet=False)
    add_sticker_label(slide, "Story > SKU", 8.2, 1.2)
    add_sticker_label(slide, "Trust Gap", 9.0, 2.0, bg_color=COLOR_INDIGO)
    add_bar_chart_mock(slide, left=6.0, top=5.5, bars=[0.7, 0.45, 0.85], labels=["Awareness", "Premium Perception", "Trust Index"])

    # SLIDE 3: SOLUTION (350+ words)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide)
    add_title(slide, "✨ Emami Sattva – The Heritage Infusion Series")

    slide3 = [
        "## Concept 🌿",
        "Emami Sattva recasts ghee as a ceremonial wellness medium—A2 Bilona Ghee infused with whole botanicals that tap into India’s Ayurvedic lineage. The hero line “Taste the Wisdom of India” isn’t marketing flourish; it frames consumption as cultural continuity. Floating ingredient visuals (turmeric, garlic) defy static commodity photography, symbolizing raw purity within motion and light. In this narrative, the product becomes an artifact of care, not just a jar of fat: sourced ethically, processed ritually, and infused thoughtfully.",
        "## The Shift 🔄",
        "This is a dual shift in process and perception: From Industrial Manufacturing → To Vedic Bilona Method; From Commodity → To Heritage Art. Operationally, this demands curation and traceability. Aesthetically, it requires editorial typography (Playfair Display) and modern clarity (Work Sans), balanced by calm cream backgrounds and Indigo anchors. Strategically, it turns consumption into ritual: a spoonful taken not for calories, but for vitality. The infusion series becomes the stage for personalization—Turmeric for immunity, Garlic for heart health, Brahmi for focus—each aligned with tangible lifestyle outcomes.",
        "## Core Value Proposition 💛",
        "\"We don’t just make ghee; we preserve Prana.\" This is the center of gravity that holds experience, process, and trust together. The proposition is substantiated by site mechanics: meaningful hover states, guided narrative sections, and traceability modules that translate claims into verifiable facts. Consumers can see beyond surface aesthetics into transparent provenance—reassurance against adulteration and a stronger sense of connection to source.",
        "## Audience & Positioning 🧘",
        "The Modern Yogi persona values rooted modernity: wellness seriousness without spiritual posturing, and heritage without nostalgia kitsch. They appreciate balanced tones, clear hierarchies, and micro-interactions that feel elegant, not gaudy. Our UI is intentionally restrained: generous whitespace, focused typography, and a golden accent that reads premium. The result: a premium brand system that can scale into new botanicals and formats while maintaining coherence. Emami Sattva isn’t a re-skin; it’s a meaningful repositioning to ritual-centered wellness—a system where the interface, product, and story honor the same principle of preserving Prana."
    ]
    add_content_text(slide, slide3, font_size=18, is_bullet=False)
    if hero_img:
        add_image(slide, hero_img, left=6.2, top=1.8, width=4.2, height=2.6)
        add_caption(slide, "Hero: Taste the Wisdom of India • Floating botanicals • Ritual minimalism", 6.2, 4.6, 4.2, 0.6)
    add_flow_step(slide, "Industrial", 6.0, 5.5, gold=False)
    add_arrow(slide, 8.8, 5.8, 0.6, 0.4)
    add_flow_step(slide, "Vedic Bilona", 9.6, 5.5, gold=True)

    # SLIDE 4: DESIGN SYSTEM (350+ words)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide)
    add_title(slide, "🎨 Visual Alchemy: Design System & UI Quality")

    slide4 = [
        "## Typography & Atmosphere 🅰️",
        "Typography sets the editorial tone for premium. Headings in a serif akin to Playfair Display deliver gravity, history, and a magazine-grade presence that immediately elevates perceived value. Body in a clean sans like Work Sans ensures contemporary readability and avoids decorative noise. This typographic tension—heritage vs modernity—embodies Emami Sattva’s promise: ancient method, present-day clarity. Generous line spacing, restrained bolding, and consistent sizing produce rhythm and calm, crucial for wellness UX.",
        "## Color Psychology ✨",
        "The palette is functional and symbolic. Cream (#FDFBF7) is breath—the space where story can unfold and product can exhale. Indigo (#1E1B4B) is trust and night-sky depth, grounding the interface and connoting Ayurvedic introspection. Gold (#D97706) is warmth, prosperity, and the liquid light of ghee itself. Used as an accent rather than a wash, Gold highlights interactions, badges, and callouts without shouting. Together, the palette says premium without performing excess.",
        "## Layout & Hierarchy 🧩",
        "Hierarchy is achieved through whitespace, scale, and pacing instead of heavy borders or loud colors. Cards in the Collection section demonstrate this: clear image focus, concise title, supportive subtext, and meaningful hover states that add delight (lift and scale) while keeping utility crisp. Glassmorphism in process areas introduces contemporary layering to rustic imagery—bridging field-to-interface with translucency that feels sophisticated, not gimmicky. Sticky navigation smooths orientation while its transparency shift on scroll keeps hero visuals unobstructed.",
        "## Systemic Scalability 🔁",
        "A well-structured design system is a growth instrument. The chosen type scales across hero headlines, product detail snippets, and educational modals. Color roles remain consistent—Indigo for anchors/labels, Gold for emphasis, Cream for rest—avoiding cognitive fatigue. Components (cards, badges, toasts, CTA buttons) are built as repeatable patterns, enabling new infusions like Ashwagandha (stress), Neem (detox), or Tulsi (respiratory) to slot into the same visual grammar without dilution. The system thus becomes a brand engine: reliable, elegant, and extensible."
    ]
    add_content_text(slide, slide4, font_size=18, is_bullet=False)

    tf = slide.placeholders[1].text_frame
    create_color_swatch(slide, 6.0, 1.2, COLOR_BG, "Sattva Base\n#FDFBF7 ✨")
    create_color_swatch(slide, 7.8, 1.2, COLOR_GOLD, "Liquid Gold\n#D97706 🟡")
    create_color_swatch(slide, 9.6, 1.2, COLOR_INDIGO, "Vedic Indigo\n#1E1B4B 🌌")
    if collection_img:
        add_image(slide, collection_img, left=6.0, top=3.0, width=4.8, height=2.0)
        add_caption(slide, "Collection card UI: hierarchy, hover lift, ingredient focus", 6.0, 5.1, 4.8, 0.6)

    # SLIDE 5: UX STRATEGY (350+ words)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide)
    add_title(slide, "🧠 UX Thinking: The ‘From Soil to Soul’ Journey")

    slide5 = [
        "## Journey Architecture 🧭",
        "The UX strategy is intentionally ceremonial, guiding users along a page-level narrative that mirrors mindful practice: Attract → Educate → Desire → Trust → Action. The hero section attracts through evocative copy (“Taste the Wisdom of India”) and floating botanicals that signal purity and infusion. Education follows through heritage explanations and a process lens, grounding claims in the Vedic Bilona method. Desire is cultivated in the Collection: Turmeric Golden Ghee (Immunity), Garlic Heart Ghee (Heart Health), Brahmi Mind Ghee (Focus). Trust is fortified by Traceability: batch codes that reveal milking dates and herdsman details (Ram Kaka, Nashik). Finally, Action is simplified through obvious CTAs and friction-minimized steps.",
        "## Interaction Design 🖱️",
        "Micro-interactions are subtle and premium. Card hover states lift and scale to acknowledge attention without garish animations; ingredient imagery shifts to emphasize botanical honesty. The sticky navbar adapts transparency on scroll, keeping orientation intact but never overshadowing hero content. Add-to-cart triggers a gentle toast notification—reassurance without interrupting the journey. Forms and modals (e.g., Ayurveda Dosha education) are concise and respectful of cognitive load, serving insight without feeling preachy.",
        "## Information Architecture 🧭",
        "Navigation clusters around three intuitive pillars: Shop (Collection and product detail), Learn (Heritage, Process, Ayurveda), and Rituals (Guides and usage). This categorization mirrors user intent and prevents category drift. Each section is designed for scannability: headings and sub-headings lead with benefit language, while body copy translates process into human terms (field, herd, method, infusion). CTA placement follows reading gravity: positioned where curiosity transforms into intent, never in the way of narrative flow.",
        "## Conversion Philosophy ✅",
        "Frictionless action is the consequence of earned trust. We reduce step count, clarify pricing, and anchor credibility with badges and traceability cues. Checkout adopts a minimal mental tax approach—clean inputs, progressive disclosure, wallet integrations where feasible. The result is a user journey that respects attention, rewards curiosity, and invites commitment as the natural next step in a ritual of care. The UX doesn’t shout; it persuades through coherence and calm."
    ]
    add_content_text(slide, slide5, font_size=18, is_bullet=False)

    add_flow_step(slide, "✨ Attract (Hero)", 6.0, 1.2)
    add_arrow(slide, 8.8, 1.5, 0.6, 0.4)
    add_flow_step(slide, "📖 Educate (Heritage)", 9.6, 1.2)
    add_flow_step(slide, "🔥 Desire (Infusions)", 6.0, 2.8)
    add_arrow(slide, 8.8, 3.1, 0.6, 0.4)
    add_flow_step(slide, "🔍 Trust (Traceability)", 9.6, 2.8)
    add_flow_step(slide, "✅ Action (CTA)", 7.8, 4.4)
    add_caption(slide, "Guided scroll → clear CTAs → minimal friction checkout 🛒", 6.0, 5.6, 5.2, 0.6)

    # SLIDE 6: INNOVATION & TRANSPARENCY (350+ words)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide)
    add_title(slide, "🛡️ Radical Transparency & Creative Tech")

    slide6 = [
        "## Traceability 🔍",
        "Trust is engineered, not assumed. The Traceability module allows users to input a batch code and reveal production specifics—milking date, process notes, and the herdsman responsible (e.g., Ram Kaka, Head Herdsman, Nashik). In a market vulnerable to adulteration, this feature converts skepticism into confidence. It reframes purchasing as a dialogue: between consumer and origin, between claim and proof. The interface is intentionally simple: a clean input, a clear result, and an elegant presentation of provenance—Indigo labels, gold accents, cream background that never distracts.",
        "## Creative Engineering 💻",
        "Premium is not only aesthetic; it’s the feeling of care in motion. A custom canvas ‘ghee cursor’ produces a soft trail of golden particles that glide with the pointer, evoking fluidity and warmth. This micro-delight reinforces product identity—liquid gold—while staying non-intrusive. Transitions and easing follow a wellness tempo: smooth, patient, reassuring. Modals for Ayurveda education are lightweight, with concise explanations of Vata, Pitta, and Kapha, helping users align purchases with personal balance without turning the site into a lecture hall.",
        "## Security & Signaling 🧪",
        "Badging systems, certification marks, and provenance highlights form a semiotic layer of trust. Placement matters: these appear adjacent to key decision junctures, not scattered randomly. Copy avoids hyperbole; it favors clarity and measured tone. Data privacy cues (where relevant) and transparent sourcing statements strengthen the perceived integrity of the brand. Together, these elements communicate that Emami Sattva values both the spirit of Ayurveda and the rigor of modern accountability.",
        "## Strategic Outcome 🧭",
        "Innovation and transparency drive differentiation beyond packaging. They transform the experience into a relationship anchored by verification and delight. The brand becomes a ritual partner: it educates gently, proves honestly, and charms subtly. That triad—education, proof, delight—creates loyalty loops that sustain premium pricing and word-of-mouth advocacy. In this configuration, every technical flourish has a purpose: to keep the promise of preserving Prana credible and emotionally resonant."
    ]
    add_content_text(slide, slide6, font_size=18, is_bullet=False)
    if traceability_img:
        add_image(slide, traceability_img, left=6.0, top=1.2, width=4.8, height=2.2)
        add_caption(slide, "Traceability: Batch code → provenance • Ram Kaka (Nashik)", 6.0, 3.5, 4.8, 0.6)
    add_sticker_label(slide, "Trust Badge ✅", 9.6, 1.0)
    add_bar_chart_mock(slide, left=6.0, top=4.4, bars=[0.4, 0.7, 0.9], labels=["Baseline Trust", "With Traceability", "With Traceability + UX Delight"])

    # SLIDE 7: BUSINESS IMPACT (350+ words)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide)
    add_title(slide, "🚀 Why This Rebrand Wins 🏆")

    slide7 = [
        "## Brand Alignment 🤝",
        "Emami’s enduring promise—Ayurvedic science—finds its most compelling digital articulation in Sattva’s ritual-centric system. By fusing heritage storytelling, traceability, and refined UI, the brand shifts from a perception of mass utility to premium wellness authority. The tonal discipline (cream, indigo, gold), editorial typography, and restrained interactions create a global-ready aesthetic that stands alongside luxury wellness brands while retaining Indian authenticity.",
        "## Differentiation ✨",
        "Commodity players compete on price and generic purity claims; Emami Sattva competes on ceremony, provenance, and benefit-led infusions. Nutrient infusions translate Ayurveda into modern outcomes—immunity (Turmeric), heart health (Garlic), focus (Brahmi)—positioning the product as both culinary and therapeutic. The experience avoids nostalgia clichés (grandma’s jar) and adopts contemporary wellness semiotics (clean layouts, calm motion, scientific cues). The ghee cursor, glassmorphism, and traceability together form a unique signature.",
        "## Scalability 📈",
        "The design system is modular and expansion-ready: Ashwagandha for stress resilience, Neem for detox, Tulsi for respiratory balance, and beyond. Each new infusion inherits a consistent pattern library (cards, badges, modals) and palette roles, guaranteeing speed to market and coherence. Educational content scales via the Learn and Rituals pillars, deepening brand authority while feeding SEO and community engagement. Operationally, traceability can extend with geo-tagging, batch histories, and farmer spotlights to enrich storytelling.",
        "## Commercial Outlook 💹",
        "Premium pricing is supported by trust mechanisms and the felt difference of ritual UX. Conversion benefits from reduced friction and elevated desirability; retention is driven by education plus delight loops. Partnerships (yoga studios, wellness clinics, mindful cooking creators) can amplify credibility, while D2C data informs personalization (dosha tips, preferred infusions, replenishment cycles). Sattva can credibly occupy the nexus of heritage and bio-hacking, appealing to global consumers seeking tradition with modern validation.",
        "## Final Thought 🌿",
        "Emami Sattva doesn’t merely sell ghee; it curates a lifestyle of purity and presence. By preserving Prana through product integrity and digital grace, the brand earns trust, commands attention, and sustains momentum. In Round 1, that coherence—between what we say, what we show, and how we behave—is the differentiator that wins."
    ]
    add_content_text(slide, slide7, font_size=18, is_bullet=False)
    add_bar_chart_mock(slide, left=6.0, top=5.5, bars=[0.5, 0.75, 0.85], labels=["Baseline", "Post-UX", "With Traceability + Infusions"])
    add_sticker_label(slide, "Heritage + Tech", 9.6, 1.1, bg_color=COLOR_INDIGO)

    # SLIDE 8: END
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    set_slide_background(slide)
    add_title(slide, "Thank You 🙏")

    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(2.5))
    tf = tb.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Q&A"
    p.font.name = "Times New Roman"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_INDIGO
    p.alignment = PP_ALIGN.LEFT

    p = tf.add_paragraph()
    p.text = "Credits: Site (cadence-rebranding-wars.netlify.app), Images (Unsplash), Icons (Lucide), Strategy (Emami Corporate Report)."
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_GOLD
    p.alignment = PP_ALIGN.LEFT

    prs.save(output_path)
    return os.path.abspath(output_path)

# --- RUN ---
if __name__ == "__main__":
    path = generate_presentation(
        output_path="Emami_Sattva_Rebrand_Round1_Expanded.pptx",
        team_name="Cadence",
        team_members="Mohit, Member 2, Member 3",
        # Provide local paths if you export screenshots from your site:
        hero_img=None,            # e.g., "assets/hero.png"
        collection_img=None,      # e.g., "assets/collection.png"
        traceability_img=None     # e.g., "assets/traceability_ram_kaka.png"
    )
    print(f"Presentation generated: {path}")