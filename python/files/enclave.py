import sys
import datetime
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# --- Define a Custom Color Palette ---
class Colors:
    DARK_GREEN = RGBColor(0x00, 0x5A, 0x31) # Deeper, professional green
    MEDIUM_GREEN = RGBColor(0x50, 0xC8, 0x78) # Emerald green for accents
    LIGHT_GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5) # Light grey background
    DARK_GRAY_TEXT = RGBColor(0x33, 0x33, 0x33) # Dark text color
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    CRISIS_RED = RGBColor(0xD9, 0x53, 0x4F)     # A softer, professional red
    SOLUTION_BLUE = RGBColor(0x02, 0x75, 0xD8) # Strong, corporate blue
    IPO_GOLD = RGBColor(0xFF, 0xA7, 0x00)     # Gold for "Advantage"

def set_run_style(run, size_pt, bold=False, color=Colors.DARK_GRAY_TEXT):
    """Helper function to style a text run."""
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(size_pt)
    font.bold = bold
    font.color.rgb = color

# --- Helper to style table cells (Robust) ---
def style_cell(cell, text, fill_color, font_color, font_size, bold=False, align=PP_ALIGN.CENTER):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    tf = cell.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    
    run = p.add_run()
    run.text = text
    set_run_style(run, font_size, bold=bold, color=font_color)

def create_presentation():
    """Main function to create the strategic Grainwave Foods PPT."""
    try:
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        # Set a light background for all slides
        slide_master = prs.slide_masters[0]
        background = slide_master.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = Colors.LIGHT_GRAY_BG
        
        print("Presentation object created. Starting Slide 1...")

        # --- Slide 1: Title Slide (No Change) ---
        slide_layout = prs.slide_layouts[0] # Title slide
        slide = prs.slides.add_slide(slide_layout)
        
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = Colors.WHITE
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Grainwave at the Crossroads: A Winning Strategic Framework"
        p = title.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run_style(p.runs[0], 44, bold=True, color=Colors.DARK_GREEN)
        
        subtitle.text = "Solving the Core Conflict Between Ethical Commitments and Profitability"
        p = subtitle.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run_style(p.runs[0], 24, color=Colors.DARK_GRAY_TEXT)
        
        print("Slide 1 (Title) created.")

        # --- Slide 2: The Core Conflict & 3-Pillar Solution (NEW SLIDE 2) ---
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "The Core Conflict & Our 3-Pillar Solution"
        set_run_style(title.text_frame.paragraphs[0].runs[0], 36, bold=True, color=Colors.DARK_GREEN)

        def add_pillar_box(text_list, left, top, color, title_text):
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4.5), Inches(5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = Colors.WHITE
            shape.line.color.rgb = color
            shape.line.width = Pt(3)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.word_wrap = True

            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title_text
            set_run_style(run, 22, bold=True, color=color)

            for item in text_list:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = item
                set_run_style(run, 18, bold=False)
                p.level = 1
            return shape

        # The Pillars
        p1 = add_pillar_box(
            ["Brand Promise (Ethics)", "vs.", "Operational Reality (Profit)", "vs.", "CFO's Margin Pressure"],
            Inches(0.5), Inches(2.0), Colors.CRISIS_RED, "The Core Conflict"
        )
        p2 = add_pillar_box(
            ["1. Hybrid Sourcing Model (Operations)", "2. IPO-Ready Governance (Strategy)", "3. 'Fair Partnership' Finance (Marketing)"],
            Inches(5.75), Inches(2.0), Colors.SOLUTION_BLUE, "The 3-Pillar Strategy"
        )
        p3 = add_pillar_box(
            ["Turn a liability into a marketable, competitive advantage.", "Satisfy all stakeholders (CFO, CSR, Investors)."],
            Inches(11.0), Inches(2.0), Colors.IPO_GOLD, "The Winning Outcome"
        )
        
        # Connectors
        try:
            # Use MSO_CONNECTOR if available; older/newer pptx enums differ
            slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, p1.left + Inches(4.5), p1.top + Inches(2.5), p2.left, p2.top + Inches(2.5))
            slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, p2.left + Inches(4.5), p2.top + Inches(2.5), p3.left, p3.top + Inches(2.5))
        except Exception:
            # Fallback: try the MSO_SHAPE constant if present, else skip connectors
            try:
                slide.shapes.add_connector(MSO_SHAPE.CONNECTOR_ELBOW, p1.left + Inches(4.5), p1.top + Inches(2.5), p2.left, p2.top + Inches(2.5))
                slide.shapes.add_connector(MSO_SHAPE.CONNECTOR_ELBOW, p2.left + Inches(4.5), p2.top + Inches(2.5), p3.left, p3.top + Inches(2.5))
            except Exception:
                pass

        print("Slide 2 (3-Pillar Strategy) created.")

        # --- Slide 3: Pillars 1 & 2: The New Operational Model (NEW SLIDE 3) ---
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Strategy Deep Dive: Sourcing & Governance"
        set_run_style(title.text_frame.paragraphs[0].runs[0], 36, bold=True, color=Colors.DARK_GREEN)

        # --- Left Side: Pie Chart (Sourcing) ---
        chart_data = CategoryChartData()
        chart_data.categories = ["Tier 1: 'Farm2Future 2.0'", "Tier 2: Open Mandi", "Tier 3: Imports/New Regions"]
        chart_data.add_series('Sourcing', (50, 30, 20))

        x, y, cx, cy = Inches(0.5), Inches(2), Inches(6), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.show_percentage = True
        chart.plots[0].data_labels.font.size = Pt(14)
        
        # Text for Pie Chart
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(6), Inches(2))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        run = p.add_run(); run.text = "Tier 1 (50%): "; set_run_style(run, 18, bold=True, color=Colors.DARK_GREEN)
        run = p.add_run(); run.text = "Secures ethical promise."; set_run_style(run, 18)
        p = tf.add_paragraph()
        run = p.add_run(); run.text = "Tier 2 (30%): "; set_run_style(run, 18, bold=True, color=Colors.SOLUTION_BLUE)
        run = p.add_run(); run.text = "Market flexibility for demand spikes."; set_run_style(run, 18)
        p = tf.add_paragraph()
        run = p.add_run(); run.text = "Tier 3 (20%): "; set_run_style(run, 18, bold=True, color=Colors.DARK_GRAY_TEXT)
        run = p.add_run(); run.text = "Long-term diversification & risk mgt."; set_run_style(run, 18)

        # --- Right Side: Text Box (Governance) ---
        txBox = slide.shapes.add_textbox(Inches(8), Inches(2), Inches(7.5), Inches(6))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Pillar 2: IPO-Ready Governance"
        set_run_style(run, 28, bold=True, color=Colors.DARK_GREEN)
        
        p = tf.add_paragraph()
        run = p.add_run(); run.text = "Action: "; set_run_style(run, 20, bold=True)
        run = p.add_run(); run.text = "Establish a permanent 'Board-Level Ethics & Supplier Relations Committee.'"; set_run_style(run, 20)
        p.level = 1
        
        p = tf.add_paragraph()
        run = p.add_run(); run.text = "Why? (Internal): "; set_run_style(run, 20, bold=True)
        run = p.add_run(); run.text = "Acts as an 'emergency brake' to ensure procurement never again breaks the brand's promise."; set_run_style(run, 20)
        p.level = 1

        p = tf.add_paragraph()
        run = p.add_run(); run.text = "Why? (External): "; set_run_style(run, 20, bold=True)
        run = p.add_run(); run.text = "Provides a crucial, positive governance signal for the upcoming IPO, satisfying investors."; set_run_style(run, 20)
        p.level = 1
        
        print("Slide 3 (Sourcing & Governance) created.")

        # --- Slide 4: Pillar 3: The 'Fair Partnership Premium' (NEW SLIDE 4) ---
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Strategy Masterstroke: The 'Fair Partnership Premium'"
        set_run_style(title.text_frame.paragraphs[0].runs[0], 36, bold=True, color=Colors.DARK_GREEN)
        
        # --- Left Side: Bar Chart (Price) ---
        chart_data = CategoryChartData()
        chart_data.categories = ['Old Price', 'New "Fair Partnership" Price']
        chart_data.add_series('Price (Rs)', (35, 38))
        
        x, y, cx, cy = Inches(0.5), Inches(2.5), Inches(6), Inches(5)
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart
        chart.has_legend = False
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.minimum_scale = 30
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.font.size = Pt(16)
        
        # Color bars
        chart.plots[0].series[0].points[0].format.fill.solid()
        chart.plots[0].series[0].points[0].format.fill.fore_color.rgb = Colors.CRISIS_RED
        chart.plots[0].series[0].points[1].format.fill.solid()
        chart.plots[0].series[0].points[1].format.fill.fore_color.rgb = Colors.MEDIUM_GREEN
        
        # --- Right Side: Text Box (The Strategy) ---
        txBox = slide.shapes.add_textbox(Inches(8), Inches(2), Inches(7.5), Inches(6))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        run = p.add_run(); run.text = "The Problem: "; set_run_style(run, 22, bold=True)
        run = p.add_run(); run.text = "New model creates a 6.4% margin hit, which is not sustainable."; set_run_style(run, 22)
        
        p = tf.add_paragraph()
        run = p.add_run(); run.text = "The Strategy: "; set_run_style(run, 22, bold=True)
        run = p.add_run(); run.text = "Turn the cost into a marketable story."; set_run_style(run, 22)
        
        p = tf.add_paragraph()
        run = p.add_run(); run.text = "The Marketing Narrative:"; set_run_style(run, 22, bold=True, color=Colors.DARK_GREEN)
        p.level = 1
        
        p = tf.add_paragraph()
        run = p.add_run(); run.text = ("\"We are proud to be transparent. 100% of this "
                                      "price increase goes directly to supporting our new, "
                                      "equitable model with our farmer partners.\"")
        set_run_style(run, 20)
        p.level = 2
        
        p = tf.add_paragraph()
        run = p.add_run(); run.text = "The Result:"; set_run_style(run, 22, bold=True, color=Colors.IPO_GOLD)
        p.level = 1
        
        p = tf.add_paragraph()
        run = p.add_run(); run.text = ("Turns a brand-destroying liability into a "
                                      "marketable, competitive advantage.")
        set_run_style(run, 20)
        p.level = 2
        
        print("Slide 4 (Finance Strategy) created.")
        
        # --- Slide 5: Summary Table (NEW SLIDE 5) ---
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Conclusion: From Liability to Lasting Advantage"
        set_run_style(title.text_frame.paragraphs[0].runs[0], 36, bold=True, color=Colors.DARK_GREEN)
        
        # Add Summary Table
        rows, cols = 4, 3
        x, y, cx, cy = Inches(1.5), Inches(2), Inches(13), Inches(5.5)
        table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
        
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(4.5)
        table.columns[2].width = Inches(4.5)
        
        # Headers
        style_cell(table.cell(0, 0), "THE CRISIS", Colors.CRISIS_RED, Colors.WHITE, 22, bold=True)
        style_cell(table.cell(0, 1), "THE STRATEGY", Colors.SOLUTION_BLUE, Colors.WHITE, 22, bold=True)
        style_cell(table.cell(0, 2), "THE ADVANTAGE", Colors.IPO_GOLD, Colors.WHITE, 22, bold=True)
        
        # Row 1
        style_cell(table.cell(1, 0), "Broken trust with farmers; brand promise destroyed.", Colors.WHITE, Colors.DARK_GRAY_TEXT, 18, align=PP_ALIGN.LEFT)
        style_cell(table.cell(1, 1), "Hybrid 50/30/20 Sourcing Model.", Colors.WHITE, Colors.DARK_GRAY_TEXT, 18, align=PP_ALIGN.LEFT)
        style_cell(table.cell(1, 2), "A resilient & ethical supply chain customers trust.", Colors.WHITE, Colors.DARK_GRAY_TEXT, 18, align=PP_ALIGN.LEFT)
        
        # Row 2
        style_cell(table.cell(2, 0), "Misaligned operations (Procurement vs. CSR).", Colors.WHITE, Colors.DARK_GRAY_TEXT, 18, align=PP_ALIGN.LEFT)
        style_cell(table.cell(2, 1), "Board-Level Ethics Committee.", Colors.WHITE, Colors.DARK_GRAY_TEXT, 18, align=PP_ALIGN.LEFT)
        style_cell(table.cell(2, 2), "Strong governance signal for a successful IPO.", Colors.WHITE, Colors.DARK_GRAY_TEXT, 18, align=PP_ALIGN.LEFT)

        # Row 3
        style_cell(table.cell(3, 0), "Unsustainable 6.4% margin hit.", Colors.WHITE, Colors.DARK_GRAY_TEXT, 18, align=PP_ALIGN.LEFT)
        style_cell(table.cell(3, 1), "'Fair Partnership Premium' (Price & Marketing).", Colors.WHITE, Colors.DARK_GRAY_TEXT, 18, align=PP_ALIGN.LEFT)
        style_cell(table.cell(3, 2), "Turns a cost into a marketable story & competitive moat.", Colors.WHITE, Colors.DARK_GRAY_TEXT, 18, align=PP_ALIGN.LEFT)
        
        print("Slide 5 (Summary Table) created.")

        # --- Save Presentation ---
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        file_path = f"Grainwave_Strategic_Solution_{timestamp}.pptx"
        
        prs.save(file_path)
        print(f"\nSuccess! The file is made and saved with this name: \"{file_path}\"")
        return file_path

    except ImportError:
        print("\n--- ERROR ---")
        print("Module 'python-pptx' not found.")
        print("Please install it by running: pip install python-pptx")
        return None
    except Exception as e:
        print(f"\nAn error occurred: {e}")
        print(f"Error details: {e}")
        import traceback
        traceback.print_exc()
        print("Please ensure you have permissions to write files in this directory.")
        return None

# --- Main execution ---
if __name__ == "__main__":
    create_presentation()
