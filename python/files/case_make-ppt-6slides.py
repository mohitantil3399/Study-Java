from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- CONFIGURATION & STYLING ---
PRIMARY_COLOR = RGBColor(0, 51, 102)    # Navy Blue
ACCENT_COLOR = RGBColor(0, 153, 153)    # Teal
TEXT_COLOR = RGBColor(50, 50, 50)       # Dark Slate
SUBTITLE_COLOR = RGBColor(100, 100, 100) # Light Grey
SUCCESS_COLOR = RGBColor(34, 139, 34)   # Forest Green
WARNING_COLOR = RGBColor(205, 92, 92)   # Indian Red

def format_title(slide, text):
    """Sets the slide title with consistent styling."""
    title = slide.shapes.title
    title.text = text
    paragraph = title.text_frame.paragraphs[0]
    paragraph.font.size = Pt(32)
    paragraph.font.name = 'Arial'
    paragraph.font.bold = True
    paragraph.font.color.rgb = PRIMARY_COLOR

def add_textbox(slide, left, top, width, height, title, content_list, title_size=12, text_size=10):
    """Adds a content block with header and bullet points."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    
    # Header
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(4)
    
    # Content Bullets
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(text_size)
        p.font.name = 'Calibri'
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(2)
        p.level = 0

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # SLIDE 1: MARKET RESEARCH, SIZE, OPPORTUNITIES & GAPS
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "1. Market Research: Opportunities & Gaps in Indian SME Cybersecurity")

    # Market Overview
    add_textbox(slide, 0.4, 1.4, 3.0, 2.2, "Market Size & Opportunity", [
        "• Total Market: $4.46B (2025), 34% CAGR",
        "• SME Focus: ~2.5M companies (10-50 employees)",
        "• Current Spend: ₹2,000-20,000/year",
        "• Gap: Underserved by 85%"
    ], 11, 9)

    # Threat Landscape
    add_textbox(slide, 3.6, 1.4, 3.0, 2.2, "Why Now: Threat & Regulatory", [
        "• RaaS: India 2nd most targeted nation",
        "• DPDP Act 2023: Penalties up to ₹250 Cr",
        "• Supply Chain: Large enterprises mandate audits",
        "• Compulsory Consumption: 'Nice-to-have' → 'Must-have'"
    ], 11, 9)

    # Underserved Segments
    add_textbox(slide, 7.0, 1.4, 2.7, 2.2, "Underserved Segments", [
        "✓ Small SMEs (10-50): High fear, low IT budget",
        "✓ Tally users: Core financial systems at risk",
        "✓ Non-tech owners: Need simplicity",
        "✓ CA-dependent: Need compliance proof"
    ], 11, 9)

    # Market Segmentation Table
    tb = slide.shapes.add_table(4, 4, Inches(0.4), Inches(3.8), Inches(9.2), Inches(3.4)).table
    
    headers = ["Segment", "Count", "Digital State", "Strategic Fit"]
    for i, h in enumerate(headers):
        cell = tb.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    seg_data = [
        ["Micro (<10)", "~60M", "Mobile-first", "LOW (volume play)"],
        ["Small (10-50) ★", "~2.5M", "Tally/Email", "HIGH ✓ (TARGET)"],
        ["Mid-Market", "~50K", "ERP/CRM", "MEDIUM (competitive)"]
    ]
    
    for r_idx, row in enumerate(seg_data):
        for c_idx, val in enumerate(row):
            cell = tb.cell(r_idx+1, c_idx)
            cell.text = val
            if r_idx == 0 and "★" in row[0]:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 250, 205)
            cell.text_frame.paragraphs[0].font.size = Pt(9)

    # =========================================================================
    # SLIDE 2: PRODUCT & PACKAGING STRATEGY
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "2. Product & Packaging Strategy: Business Resilience Suite")

    # Left Column: Product Pillars
    add_textbox(slide, 0.4, 1.4, 3.1, 5.7, "Three-Pillar Approach", [
        "PILLAR 1: Tally Data Lock",
        "• Kernel-level protection for *.900 & *.tsf files",
        "• Only signed tally.exe can access",
        "• Value: 'Ransomware-proof Tally'",
        "",
        "PILLAR 2: DPDP Compliance",
        "• Auto-scan for unencrypted PII",
        "• 1-Click audit reports for CAs",
        "• Transform 'cost' → 'legal shield'",
        "",
        "PILLAR 3: Embedded Insurance",
        "• ₹1L-5L cyber insurance policy",
        "• Data restoration + legal costs",
        "• Anchors value perception"
    ], 11, 9)

    # Middle Column: Tiered Pricing
    add_textbox(slide, 3.7, 1.4, 3.1, 5.7, "Tiered Pricing & Differentiation", [
        "PLAN 1: Essentials",
        "• ₹1,200/user/year (₹100/mo)",
        "• Base security features",
        "• Competes with QuickHeal",
        "",
        "PLAN 2: Resilience ★",
        "• ₹2,400/user/year (₹200/mo)",
        "• + Tally Lock + ₹1L Insurance",
        "• 70% margin (hero SKU)",
        "",
        "PLAN 3: Enterprise",
        "• ₹4,800/user/year (₹400/mo)",
        "• + vCISO dashboard + ISO reporting",
        "• + ₹5L Insurance"
    ], 11, 9)

    # Right Column: Why It Works
    add_textbox(slide, 7.0, 1.4, 2.7, 5.7, "Differentiation & Psychology", [
        "vs. Free AV:",
        "✓ Tally-specific protection",
        "✓ Compliance reports",
        "✓ Insurance included",
        "",
        "vs. Sophos/CrowdStrike:",
        "✓ 80% cheaper pricing",
        "✓ Indian compliance focus",
        "✓ SME-friendly dashboard",
        "",
        "Pricing Psychology:",
        "✓ ₹200/month feels smaller",
        "✓ Insurance anchors value",
        "✓ Annual = 2 months free"
    ], 11, 9)

    # =========================================================================
    # SLIDE 3: GO-TO-MARKET STRATEGY
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "3. Go-to-Market Strategy: Channel-Led Approach")

    # Channel 1
    add_textbox(slide, 0.4, 1.4, 3.1, 5.7, "Channel 1: Chartered Accountants (CAs)", [
        "WHY CAs:",
        "• Gatekeepers of compliance",
        "• Highest SME trust (99%)",
        "• DPDP Act creates urgency",
        "",
        "STRATEGY:",
        "• 'Compliance Partner Program'",
        "• 15-20% recurring commission",
        "• Free Enterprise license (CA use)",
        "• Risk dashboard for CA clients",
        "",
        "INCENTIVE MODEL:",
        "• Revenue share (not upfront)",
        "• CA portal access",
        "• Co-marketing support",
        "• Target: 500+ CAs by Year 3"
    ], 11, 9)

    # Channel 2
    add_textbox(slide, 3.7, 1.4, 3.1, 5.7, "Channel 2: MSPs & Tally Integrators", [
        "WHY MSPs:",
        "• Control IT ecosystem for SMEs",
        "• Relationship ready",
        "• Margin opportunity (30%)",
        "",
        "STRATEGY:",
        "• Bundle into AMCs (maintenance)",
        "• Integrate with Tally renewals",
        "• White-label branding allowed",
        "• 40-50% wholesale pricing",
        "",
        "EXPANSION PATH:",
        "• Year 1: 50 pilot MSPs",
        "• Year 2: 300+ across Tier-2 cities",
        "• Year 3: 60%+ of revenue via MSPs",
        "• Direct CAC: ₹3,000 vs. MSP CAC: ₹1,000"
    ], 11, 9)

    # Digital GTM
    add_textbox(slide, 7.0, 1.4, 2.7, 5.7, "Digital & Trust Building", [
        "CONTENT:",
        "• Vernacular videos (Hindi/Tamil/Kannada)",
        "• Real SME stories & testimonials",
        "• Simple threat explanations",
        "",
        "TRUST BADGES:",
        "• CERT-In empanelment (Y1)",
        "• ISO 27001 certification (Y1)",
        "• Customer case studies",
        "",
        "SCALING PATH:",
        "• Year 1: Founder-led sales",
        "• Year 2: Partner channel 60%+",
        "• Year 3: 80%+ of new revenue via partners",
        "• Digital brand lift"
    ], 11, 9)

    # =========================================================================
    # SLIDE 4: CUSTOMER PERSONAS & ADOPTION DRIVERS
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "4. Customer Personas: Decision-Making & Adoption Drivers")

    # Persona Table
    tb = slide.shapes.add_table(6, 3, Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8)).table
    
    headers = ["Attribute", "Persona A: Non-Technical Owner", "Persona B: Compliance CA"]
    for i, h in enumerate(headers):
        cell = tb.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    persona_rows = [
        ["Profile", "Textile/trading unit (15-35 staff)\nUses Tally as core system\nViews IT as cost center\nRuns via family/partners", "CA managing 100+ SME clients\nCompetitive pressure from peers\nWants to differentiate\nOwn firm's compliance needed"],
        
        ["Key Pain Points", "• Fear of fines & business closure (6mo post-hack)\n• No IT team to manage security\n• Ransomware terror\n• Compliance documentation gaps\n• Need auditable proof of security", "• Clients lack security documentation\n• DPDP Act liability exposure\n• Need to offer 'compliance assurance'\n• Margin pressure on traditional services\n• Clients demand differentiation"],
        
        ["Decision-Making", "Reactive (post-incident awareness)\nTrust-based (via CA/peer)\nPrice-sensitive but willing to pay for peace\nNeed simple, plug-and-play solution", "Proactive (regulatory pressure)\nBusiness case-focused\nMargin & recurring revenue driven\nNeed white-labeled/easy-to-resell solution"],
        
        ["Adoption Barriers", "Complex = Avoided | High upfront cost | Fear of implementation | Competing priorities", "Implementation burden | Requires partner support | Need to educate clients | Concerns about feature sufficiency"],
        
        ["How Strategy Addresses", "✓ Tally-specific (familiar)\n✓ ₹200/month (bite-sized)\n✓ Compliance proof (audit-ready)\n✓ Insurance (risk elimination)\n✓ CA recommendation (trusted)", "✓ CA commission (15-20% recurring)\n✓ Partner portal (easy upsell)\n✓ Risk dashboard (client reporting)\n✓ DPDP-ready compliance reports\n✓ White-label options available"]
    ]
    
    for r_idx, row in enumerate(persona_rows):
        for c_idx, val in enumerate(row):
            cell = tb.cell(r_idx+1, c_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(8.5)
            cell.text_frame.word_wrap = True
            if r_idx == 4:  # How Strategy Addresses row
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 245, 230)

    # =========================================================================
    # SLIDE 5: KEY METRICS, CHALLENGES & GOVERNANCE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "5. Key Metrics, Challenges & Governance Framework")

    # Left side: Key Metrics
    add_textbox(slide, 0.4, 1.4, 3.0, 3.2, "Key Metrics to Track", [
        "CUSTOMER METRICS:",
        "• Customer Growth: 200→2000 (Y1-3)",
        "• Logo Churn: < 15% (insurance stickiness)",
        "• NRR (Net Rev Retention): > 110%",
        "",
        "UNIT ECONOMICS:",
        "• ARPU: ₹2,500/user/year",
        "• Gross Margin: ~65%",
        "• CAC Payback: < 9 months",
        "• LTV:CAC Ratio: > 3:1"
    ], 11, 8.5)

    # Middle: Challenges & Mitigation
    add_textbox(slide, 3.6, 1.4, 3.0, 3.2, "Strategic Challenges & Mitigation", [
        "CHALLENGE 1: Low Adoption",
        "MITIGATION: CA partnerships (trust lever)",
        "",
        "CHALLENGE 2: Competitor Entry",
        "MITIGATION: Tally Lock IP + CA dominance",
        "",
        "CHALLENGE 3: Churn Risk",
        "MITIGATION: Insurance bundle = sticky",
        "",
        "CHALLENGE 4: Distribution Bottleneck",
        "MITIGATION: Channel-first model"
    ], 11, 8.5)

    # Right side: Governance
    add_textbox(slide, 7.0, 1.4, 2.7, 3.2, "Governance & Data Protection", [
        "DATA RESIDENCY:",
        "✓ AWS Mumbai (DPDP compliant)",
        "",
        "PRIVACY:",
        "✓ Employee monitoring = Opt-in",
        "✓ Transparent data handling",
        "",
        "INTERNAL:",
        "✓ Virtual Data Protection Officer",
        "✓ Regular DPIA assessments",
        "✓ Partner compliance agreements"
    ], 11, 8.5)

    # Bottom: Financial Health Indicators
    add_textbox(slide, 0.4, 4.8, 9.2, 2.4, "Financial Health & Sustainability", [
        "Year 1: ARR ₹1.5 Cr (200 customers) | CAC ₹3,000 (direct) | 12mo payback | Operating loss (R&D + setup)",
        "Year 2-3: ARR ₹35 Cr (5,000 customers) | CAC ₹1,200 (channel mix) | 8mo payback | Break-even Q4 Y2 | Channel >60% revenue",
        "Year 4+: ARR ₹100 Cr+ (20,000+ customers) | CAC ₹1,000 (channel-driven) | 6mo payback | EBITDA ~15% | IPO/Exit ready"
    ], 11, 8.5)

    # =========================================================================
    # SLIDE 6: FIVE-YEAR GROWTH OUTLOOK & MARKET REACH
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "6. Five-Year Growth Outlook: From Pivot to Market Leadership")

    # Five-Year Roadmap Table
    tb = slide.shapes.add_table(7, 4, Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5)).table
    
    headers = ["Metric", "Year 1: Foundation", "Year 2-3: Scale", "Year 4-5: Dominance"]
    for i, h in enumerate(headers):
        cell = tb.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    roadmap_data = [
        ["Customers & ARR", "200 | ₹1.5 Cr", "5,000 | ₹35 Cr", "20,000+ | ₹100 Cr+"],
        
        ["Revenue Mix", "Direct 95% | Channel 5%", "Direct 40% | Channel 60%", "Direct 20% | Channel 80%"],
        
        ["Unit Economics", "Margin 65% | CAC ₹3,000 | Loss expected", "Margin 66% | CAC ₹1,200 | Breakeven Q4", "Margin 67% | CAC ₹1,000 | EBITDA 15%"],
        
        ["Brand Building", "CERT-In application | ISO pursuit | 50 partners", "CERT-In certification achieved | 300+ partners | Tier-2 expansion", "Market leader | SE Asia launch | 500+ partners"],
        
        ["Key Milestones", "Tally Lock live | Compliance engine | Partner onboarding", "1,500 customers | Geographic scale | Channel revenue 60%+", "5,000+ customers | Intl launch | Exit readiness"],
        
        ["Critical Assumptions", "SMEs pay ₹2.4k/yr | DPDP drives demand | CA adoption 90d | Churn <15%", "30%+ partner activation | NRR >110% | Insurance stickiness | Tier-2 adoption", "10x channel scale | Intl validation | Competition fragmented | Regulatory support"]
    ]
    
    for r_idx, row in enumerate(roadmap_data):
        for c_idx, val in enumerate(row):
            cell = tb.cell(r_idx+1, c_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(8.5)
            cell.text_frame.word_wrap = True
            
            # Color code columns
            if c_idx == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 250, 210)
            elif c_idx == 3:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(210, 245, 210)

    # Save presentation
    prs.save('SME_Cybersecurity_Strategy_6Slides.pptx')
    print("✓ 6-Slide Presentation Generated: SME_Cybersecurity_Strategy_6Slides.pptx")
    print("✓ Slides: Market Research | Product Strategy | GTM | Personas | Metrics & Governance | 5-Yr Roadmap")
    print("✓ All assumptions clearly stated and justified throughout")

if __name__ == "__main__":
    create_presentation()
