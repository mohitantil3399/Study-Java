from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- CONFIGURATION & STYLING ---
PRIMARY_COLOR = RGBColor(0, 51, 102)    # Navy Blue
ACCENT_COLOR = RGBColor(0, 153, 153)    # Teal
TEXT_COLOR = RGBColor(50, 50, 50)       # Dark Slate
SUBTITLE_COLOR = RGBColor(100, 100, 100) # Light Grey
SUCCESS_COLOR = RGBColor(34, 139, 34)   # Forest Green
WARNING_COLOR = RGBColor(205, 92, 92)   # Indian Red
INFO_COLOR = RGBColor(70, 130, 180)     # Steel Blue

def format_title(slide, text):
    """Sets the slide title with consistent styling."""
    title = slide.shapes.title
    title.text = text
    paragraph = title.text_frame.paragraphs[0]
    paragraph.font.size = Pt(32)
    paragraph.font.name = 'Arial'
    paragraph.font.bold = True
    paragraph.font.color.rgb = PRIMARY_COLOR

def add_textbox(slide, left, top, width, height, title, content_list, title_size=13, text_size=10):
    """Adds a content block with header and bullet points."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    
    # Header
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(4)
    
    # Content Bullets
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(text_size)
        p.font.name = 'Calibri'
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(3)
        p.level = 0

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # SLIDE 1: EXECUTIVE SUMMARY - Current State
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "1. Executive Summary: The Strategic Diagnosis")
    
    add_detailed_textbox(slide, 0.5, 1.5, 4.7, 3.2, "Current State Analysis", [
        "• ARR: ₹265 Lakhs from 18 customers",
        "• Current pricing: ₹36,000/year (friction point)",
        "• Sales model: Founder-led (unscalable)",
        "• Market perception: 'Expensive & Complex'",
        "• Product-Market Fit: Partially achieved",
        "• Key Issue: Positioning mismatch"
    ])
    
    add_detailed_textbox(slide, 5.3, 1.5, 4.2, 3.2, "Three Critical Failures", [
        "1. MISALIGNED VALUE: Selling 'prevention' vs. buying 'business continuity'",
        "2. FRICTION PRICING: ₹36k upfront without clear ROI is barrier",
        "3. NON-SCALABLE GTM: Founder-led sales lacks trust levers needed for scale"
    ])
    
    add_detailed_textbox(slide, 0.5, 4.8, 9.2, 2.4, "The Strategic Pivot: 'Resilience Over Security'", [
        "Transform from 'Cybersecurity Vendor' → 'Business Resilience Partner'",
        "Bundle: Endpoint Protection + Automated Compliance (DPDP) + Embedded Insurance",
        "Distribution: Direct → Channel-First (CAs & MSPs)",
        "5-Year Target: ₹100Cr+ ARR, 20,000+ customers, IPO/Acquisition Exit"
    ])

    # =========================================================================
    # SLIDE 2: MARKET LANDSCAPE - DETAILED SEGMENTATION
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "2. Market Landscape: Deep Dive into SME Segmentation")
    
    # Detailed Segmentation Table
    tb = slide.shapes.add_table(5, 6, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.5)).table
    
    headers = ["Segment", "Count", "Digital Maturity", "Current Spend", "Propensity", "Strategic Fit"]
    for i, h in enumerate(headers):
        cell = tb.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    rows = [
        ["Micro (<10)", "~60M", "Low (Mobile)", "₹0-500/yr", "Low", "LOW"],
        ["Small (10-50)", "~2.5M", "Medium (Tally)", "₹2k-20k/yr", "HIGH", "HIGH ✓"],
        ["Mid-Market", "~50k", "High (ERP)", "₹1L-10L/yr", "High", "MEDIUM"],
        ["Focus Target", "2.5M units", "Digitized but no IT", "Seeking solution", "Fear of loss", "FOCUS HERE"]
    ]
    
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tb.cell(r_idx+1, c_idx)
            cell.text = val
            if r_idx == 3:  # Highlight focus row
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 245, 230)
            cell.text_frame.paragraphs[0].font.size = Pt(9)

    # =========================================================================
    # SLIDE 3: THREAT LANDSCAPE - WHY NOW
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "3. The Threat Landscape: Why 'Good Enough' is Failing")
    
    add_detailed_textbox(slide, 0.5, 1.5, 3.0, 3.0, "Ransomware-as-a-Service (RaaS)", [
        "• India: 2nd most targeted nation",
        "• Attack Vector: Automated tools scan for open ports in Tally",
        "• Outcome: Encryption of financial data = business closure",
        "• Timeframe: 6 months to closure post-attack"
    ])
    
    add_detailed_textbox(slide, 3.6, 1.5, 3.0, 3.0, "The 'Human' Vector", [
        "• AI-crafted WhatsApp phishing attacks",
        "• Bypass technical firewalls via social engineering",
        "• Target: Non-technical staff",
        "• Localized threats specific to India"
    ])
    
    add_detailed_textbox(slide, 7.0, 1.5, 2.7, 3.0, "Supply Chain Risk", [
        "• Large enterprises mandate vendor audits",
        "• Non-compliant SMEs lose contracts",
        "• Risk: Multi-year revenue loss",
        "• Pressure: Increasing exponentially"
    ])
    
    add_detailed_textbox(slide, 0.5, 4.7, 9.2, 2.4, "The Regulatory Catalyst: DPDP Act 2023 (Compulsory Consumption)", [
        "PENALTIES: Up to ₹250 Crore for failure to implement 'reasonable safeguards'  |  DATA LIABILITY: Every SME processing personal data is a Data Fiduciary",
        "GRIEVANCE REDRESSAL: Must respond to data principal queries  |  AUDIT REQUIREMENTS: Maintain logs & trails for breach proof"
    ])

    # =========================================================================
    # SLIDE 4: PRODUCT STRATEGY - BUSINESS RESILIENCE SUITE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "4. Product Strategy: The 'Business Resilience Suite'")

    # Intro text
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Pivot: 'Commoditized Antivirus' → 'Business Continuity Partner'"
    p.font.italic = True
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WARNING_COLOR

    # Pillar 1: Tally Lock
    add_detailed_textbox(slide, 0.5, 1.8, 3.0, 5.3, "Pillar 1: 'Tally Data Lock'", [
        "THE ANCHOR FEATURE",
        "• Kernel-level driver",
        "• 'Force Field' around *.900 & *.tsf files",
        "• Rejects unauthorized writes",
        "• Only signed tally.exe access",
        "",
        "VALUE PROP:",
        "'Only solution guaranteeing Tally immunity to ransomware'"
    ])

    # Pillar 2: Automated Compliance
    add_detailed_textbox(slide, 3.6, 1.8, 3.0, 5.3, "Pillar 2: DPDP Compliance", [
        "THE VALUE ADD",
        "• Data discovery (PAN, Aadhar, GSTIN)",
        "• Scan for unencrypted PII",
        "• 1-Click audit reports",
        "• Business Risk Score (0-100)",
        "",
        "TRANSFORM FROM:",
        "'Sunk cost' → 'Legal shield'"
    ])

    # Pillar 3: Embedded Insurance
    add_detailed_textbox(slide, 7.0, 1.8, 2.7, 5.3, "Pillar 3: Cyber Insurance", [
        "THE WARRANTY",
        "• Group policy via insurtech",
        "• ₹1L-5L cover",
        "• Data restoration + legal",
        "• Business interruption",
        "",
        "IMPACT:",
        "₹2,400 fee feels trivial vs. ₹5,000 insurance"
    ])

    # =========================================================================
    # SLIDE 5: PRICING STRATEGY - PSYCHOLOGICAL & TIERED
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "5. Pricing Strategy: From Cost-Plus to Value-Based")

    # Pricing Psychology
    add_detailed_textbox(slide, 0.5, 1.5, 4.5, 2.0, "Pricing Psychology (Indian Market)", [
        "✓ Sachet Pricing: ₹200/month > ₹2,400/year (bite-sized preference)",
        "✓ Anchoring: Insurance worth ₹5,000 makes ₹2,400 seem trivial",
        "✓ Bundling: Per-user packs with flat fee options",
        "✓ Incentives: Annual pre-payment + 2-Month Free"
    ])

    # Tiered Pricing Table
    tb = slide.shapes.add_table(4, 5, Inches(0.5), Inches(3.6), Inches(9.0), Inches(3.5)).table
    
    headers = ["Feature", "Essentials", "Business Resilience ★", "Enterprise", "Max Users"]
    for i, h in enumerate(headers):
        cell = tb.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR if i != 2 else ACCENT_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    pricing_rows = [
        ["Price/User/Year", "₹1,200", "₹2,400 ★", "₹4,800", "N/A"],
        ["Tally Lock + Compliance", "✓ Base", "✓ Full", "✓ Full", "N/A"],
        ["Cyber Insurance", "None", "₹1L Cover", "₹5L Cover", "Included"]
    ]
    
    for r_idx, row in enumerate(pricing_rows):
        for c_idx, val in enumerate(row):
            cell = tb.cell(r_idx+1, c_idx)
            cell.text = val
            if c_idx == 2:  # Highlight recommended
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 250, 205)
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            if "✓" in val or "★" in val:
                cell.text_frame.paragraphs[0].font.bold = True

    # =========================================================================
    # SLIDE 6: GO-TO-MARKET - CHANNEL STRATEGY
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "6. Go-to-Market: Channel-Led Revolution")

    # Channel 1: CAs
    add_detailed_textbox(slide, 0.5, 1.5, 3.0, 3.3, "Channel 1: Chartered Accountants", [
        "WHY CAs?",
        "• Gatekeepers of compliance",
        "• Highest SME trust",
        "",
        "STRATEGY:",
        "• DPDP Compliance pitch",
        "• 15-20% recurring commission",
        "• Free Enterprise license",
        "• Risk dashboard for clients"
    ])

    # Channel 2: MSPs
    add_detailed_textbox(slide, 3.6, 1.5, 3.0, 3.3, "Channel 2: MSPs & Integrators", [
        "WHY MSPs?",
        "• Own SME relationships",
        "• Control Tally/IT ecosystem",
        "",
        "STRATEGY:",
        "• 30% margin opportunity",
        "• Integrate into AMCs",
        "• White-label branding",
        "• 40-50% wholesale pricing"
    ])

    # Channel 3: Digital
    add_detailed_textbox(slide, 7.0, 1.5, 2.7, 3.3, "Channel 3: Digital GTM", [
        "CONTENT:",
        "• Vernacular videos",
        "• Hindi/Tamil/Kannada",
        "",
        "TRUST BADGES:",
        "• CERT-In empanelment",
        "• ISO 27001 cert",
        "• Customer testimonials"
    ])

    add_detailed_textbox(slide, 0.5, 4.9, 9.2, 2.2, "Channel Economics & Positioning", [
        "DIRECT SALES CAC: ₹3,000/customer  |  CHANNEL CAC: ₹1,000/customer (60% reduction via revenue share)",
        "TARGET: Year 1 (50 partners) → Year 3 (500+ partners with 60% revenue contribution)  |  MATURITY: Partner channel becomes 80%+ of revenue by Year 4-5"
    ])

    # =========================================================================
    # SLIDE 7: CUSTOMER PERSONAS
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "7. Customer Personas: Understanding the Buyer")

    # Persona Table
    tb = slide.shapes.add_table(5, 3, Inches(0.4), Inches(1.5), Inches(9.2), Inches(5.5)).table
    
    headers = ["Attribute", "Persona A: Non-Technical Owner", "Persona B: Compliance CA"]
    for i, h in enumerate(headers):
        cell = tb.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)

    persona_data = [
        ["Profile", "Runs textile/trading (10-50 staff)\nUses Tally\nViews IT as cost center", "Manages tax for 100+ SMEs\nCompetitive pressure\nWants to be trusted advisor"],
        ["Pain Points", "Fear of fines & business closure\n(6 months post-hack)\nNo IT team\nNeed proof of security", "Clients lack compliance docs\nNeed audit trail proof\nRisk from DPDP Act liability\nDifferentiation from competitors"],
        ["Buying Trigger", "Post-ransomware awareness\nRegulatory pressure\nTrust in CA recommendations", "DPDP Act enforcement\nClient demand\nChannel opportunity"],
        ["Value Driver", "Tally protection (top priority)\nSimple to implement\nPrice < ₹300/month\nProof of compliance", "DPDP readiness\nClient risk dashboard\nRecurring revenue\nZero implementation burden"]
    ]
    
    for r_idx, row in enumerate(persona_data):
        for c_idx, val in enumerate(row):
            cell = tb.cell(r_idx+1, c_idx)
            cell.text = val
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 248, 255)
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.word_wrap = True

    # =========================================================================
    # SLIDE 8: OPERATIONAL METRICS & KPIs
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "8. Operational Execution: KPIs & Governance")

    # KPI Table
    tb = slide.shapes.add_table(6, 4, Inches(0.4), Inches(1.5), Inches(9.2), Inches(3.2)).table
    
    kpi_headers = ["Metric", "Target", "Rationale", "Current"]
    for i, h in enumerate(kpi_headers):
        cell = tb.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    kpi_data = [
        ["CAC Payback Period", "< 9 Months", "Sustainability of cash flow", "TBD"],
        ["NRR (Net Rev Retention)", "> 110%", "Upsell > Churn indicator", "TBD"],
        ["Partner Activation", "> 30% in 90d", "Quality of partner pipeline", "TBD"],
        ["Compliance Report Gen.", "> 60% users", "Product stickiness proxy", "TBD"],
        ["Logo Churn", "< 15% annual", "SME retention via insurance", "Monitor"]
    ]
    
    for r_idx, row in enumerate(kpi_data):
        for c_idx, val in enumerate(row):
            cell = tb.cell(r_idx+1, c_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(9)

    # Governance Box
    add_detailed_textbox(slide, 0.4, 4.8, 9.2, 2.4, "Data Governance & Privacy (DPDP Compliance)", [
        "DATA RESIDENCY: AWS Mumbai region (mandatory)  |  PRIVACY BY DESIGN: Employee monitoring opt-in only",
        "PARTNER COMPLIANCE: Anti-bribery & data privacy agreements  |  INTERNAL AUDIT: Virtual DPO + regular DPIA assessments"
    ])

    # =========================================================================
    # SLIDE 9: UNIT ECONOMICS
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "9. Unit Economics: The Math Behind the Model")

    # Left side: Unit Economics
    add_detailed_textbox(slide, 0.5, 1.5, 4.5, 3.0, "Blended Unit Economics", [
        "ARPU: ₹2,500/user/year",
        "",
        "GROSS MARGIN: ~65%",
        "  - Cloud costs: ₹400",
        "  - Insurance premium: ₹300",
        "  - Support: ₹100",
        "  - Gross profit: ₹1,700 (68%)",
        "",
        "Note: Lower than pure SaaS (80%) but offset by CAC efficiency"
    ])

    # Right side: CAC Efficiency
    add_detailed_textbox(slide, 5.3, 1.5, 4.2, 3.0, "Customer Acquisition Cost (CAC)", [
        "DIRECT SALES CAC:",
        "  ₹3,000/customer",
        "  (High touch, founder-led)",
        "",
        "CHANNEL PARTNER CAC:",
        "  ₹1,000/customer",
        "  (Revenue share model)",
        "",
        "EFFICIENCY: 67% reduction via channel"
    ])

    # Bottom: LTV Calculation
    add_detailed_textbox(slide, 0.5, 4.7, 9.2, 2.3, "Customer Lifetime Value (LTV) Analysis", [
        "ASSUMPTION: 3-year average customer lifetime | CHURN: 15% annual (SME baseline)",
        "LTV = ARPU × 3 years × Gross Margin × Retention = ₹2,500 × 3 × 68% × 85% × 85% ≈ ₹3,700",
        "LTV:CAC Ratio = ₹3,700 : ₹1,000 (Channel) = 3.7:1 (Healthy; >3:1 is sustainable)"
    ])

    # =========================================================================
    # SLIDE 10: FIVE-YEAR FINANCIAL ROADMAP
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "10. Five-Year Growth Roadmap: From Pivot to Dominance")

    # Big table with phases
    tb = slide.shapes.add_table(5, 4, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.5)).table
    
    headers = ["Phase", "Year 1: Foundation", "Year 2-3: Scale", "Year 4-5: Dominance"]
    for i, h in enumerate(headers):
        cell = tb.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)

    five_year = [
        ["Customers", "200\n₹1.5 Cr ARR", "5,000\n₹35 Cr ARR", "20,000+\n₹100 Cr+ ARR"],
        ["Focus", "Product pivot\n(Tally Lock + Compliance)\nPartner recruitment (50)", "Tier-2 cities expansion\nVernacular marketing\nChannel scaling (300+)", "International (SE Asia)\nvCISO AI features\nExit readiness"],
        ["Financials", "High burn\nR&D + Setup\nOperating loss\nLoss/customer", "Break-even Q4 Y2\nUnit econ stabilize\nCAC down 50%\n>60% channel revenue", "EBITDA ~15%\nHigh profitability\nSeries A/B ready\nIPO/M&A path"],
        ["Metrics", "CAC payback: 12mo\nChurn: 20%\nNRR: 105%", "CAC payback: 8mo\nChurn: 15%\nNRR: 115%", "CAC payback: 6mo\nChurn: 10%\nNRR: 125%"]
    ]
    
    for r_idx, row in enumerate(five_year):
        for c_idx, val in enumerate(row):
            cell = tb.cell(r_idx+1, c_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.word_wrap = True
            
            # Color code by phase
            if c_idx == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 250, 210)  # Year 1 highlight
            elif c_idx == 3:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(210, 245, 210)  # Year 4-5 highlight

    # =========================================================================
    # SLIDE 11: STRATEGIC RISKS & MITIGATION
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "11. Strategic Risks & Mitigation Framework")

    # Risk 1
    add_detailed_textbox(slide, 0.5, 1.5, 3.0, 3.2, "Risk 1: Churn", [
        "PROBLEM:",
        "• High SME churn (20%+)",
        "• Commoditized product",
        "",
        "MITIGATION:",
        "• Focus 10-50 user segment",
        "• Auto-renewal incentives",
        "• Insurance bundle stickiness"
    ])

    # Risk 2
    add_detailed_textbox(slide, 3.6, 1.5, 3.0, 3.2, "Risk 2: Competition", [
        "PROBLEM:",
        "• Copycat competitors",
        "• Global vendors entering",
        "",
        "MITIGATION:",
        "• Tally Lock proprietary IP",
        "• CA channel dominance",
        "• Speed to market advantage"
    ])

    # Risk 3
    add_detailed_textbox(slide, 7.0, 1.5, 2.7, 3.2, "Risk 3: Insurance", [
        "PROBLEM:",
        "• Claims fraud risk",
        "• Partner integration",
        "",
        "MITIGATION:",
        "• Use established partners",
        "• Claims verification",
        "• Risk underwriting"
    ])

    add_detailed_textbox(slide, 0.5, 4.8, 9.2, 2.2, "Critical Success Factors (CSFs)", [
        "1. PRODUCT EXCELLENCE: Tally Lock must be bulletproof (zero false positives)  |  2. CHANNEL ACTIVATION: 30%+ of partners close deals in 90 days",
        "3. BRAND TRUST: CERT-In & ISO 27001 certifications within 12 months  |  4. FINANCIAL DISCIPLINE: Unit economics + LTV:CAC > 3:1"
    ])

    # =========================================================================
    # SLIDE 12: 90-DAY SPRINT PLAN
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    format_title(slide, "12. Implementation: 90-Day Sprint Plan")

    add_detailed_textbox(slide, 0.5, 1.5, 3.0, 3.0, "MONTH 1: Product", [
        "• Develop Tally Data Lock driver",
        "• Test with Tally Prime",
        "• Beta with 5 customers",
        "• Integration testing",
        "• Security hardening"
    ])

    add_detailed_textbox(slide, 3.6, 1.5, 3.0, 3.0, "MONTH 2: Integration", [
        "• Integrate insurtech API",
        "• Build compliance report generator",
        "• Dashboard enhancements",
        "• CA/MSP portal setup",
        "• User training materials"
    ])

    add_detailed_textbox(slide, 7.0, 1.5, 2.7, 3.0, "MONTH 3: Launch", [
        "• Beta launch 'Resilience' tier",
        "• Partner onboarding (50)",
        "• DPDP webinars (ICAI)",
        "• Marketing collateral",
        "• Sales training"
    ])

    add_detailed_textbox(slide, 0.5, 4.7, 9.2, 2.3, "Parallel Track: Go-Live Checklist", [
        "PRODUCT: Tally Lock driver, Compliance reports, Insurance integration  |  CHANNEL: CA/MSP portal, Revenue share agreements, Sales battlecards",
        "GOVERNANCE: DPO appointment, DPIA assessment, Privacy policy update  |  MARKETING: Website overhaul, Video content (3 languages), CERT-In application"
    ])

    # =========================================================================
    # SLIDE 13: CONCLUSION
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank for custom design
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 248, 255)
    background.line.color.rgb = RGBColor(240, 248, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "The Roadmap to Resilience"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Content
    conclusion_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8.0), Inches(5.5))
    tf = conclusion_box.text_frame
    tf.word_wrap = True
    
    content = [
        ("1. THE CORE PROBLEM", "Indian SMEs don't buy 'Security'; they buy 'Trust' and 'Business Continuity'."),
        ("2. THE PIVOT", "Transform into 'Business Resilience Partner': Tally Protection + DPDP Compliance + Embedded Insurance."),
        ("3. THE EXECUTION", "Channel-First strategy (CAs & MSPs) solves distribution. Sachet pricing removes friction."),
        ("4. THE OUTCOME", "This roadmap transforms regulatory headwinds into tailwinds, defining the 'SME Cyber Resilience' category."),
    ]
    
    for idx, (heading, text) in enumerate(content):
        p = tf.add_paragraph()
        p.text = heading
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)
        p.level = 0
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.6))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Target: ₹100 Cr+ ARR | 20,000+ Customers | Year 4-5 IPO/Exit"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_COLOR

    # Save presentation
    prs.save('SME_Cybersecurity_Detailed_Strategy.pptx')
    print("✓ Enhanced Presentation Generated: SME_Cybersecurity_Detailed_Strategy.pptx")
    print("✓ Total Slides: 13 (Cover + Executive Summary + Market + Product + Pricing + GTM + Personas + KPIs + Unit Economics + 5-Yr Roadmap + Risks + 90-Day + Conclusion)")
    print("✓ Enhanced with: Detailed tables, KPI metrics, persona analysis, financial modeling, implementation plan, governance framework")


if __name__ == "__main__":
    create_presentation()