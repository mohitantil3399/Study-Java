from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def set_slide_background(slide, rgb_color):
    """Sets the slide background to a solid RGB color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False, font_color=RGBColor(255, 255, 255), align='left'):
    """Adds a styled textbox to the slide."""
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Inter'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color):
    """Adds a styled rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

# --- Color Palette (Neon Theme) ---
BG_COLOR = RGBColor(31, 41, 55)      # bg-gray-800
CARD_BG_COLOR = RGBColor(17, 24, 39) # bg-gray-900
TEXT_COLOR_LIGHT = RGBColor(209, 213, 219) # text-gray-300
NEON_ORANGE = RGBColor(249, 115, 22)
NEON_CYAN = RGBColor(6, 182, 212)
BORDER_COLOR = RGBColor(75, 85, 99) # border-gray-600

# --- Presentation Setup ---
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_slide_layout = prs.slide_layouts[6] # Blank layout
slide = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide, BG_COLOR)

# --- Header ---
add_textbox(slide, "QuickByte: Tier-2 EXPANSION", Inches(0.5), Inches(0.2), Inches(15), Inches(1), font_size=44, bold=True, font_color=NEON_ORANGE)
add_textbox(slide, "A High-Voltage Playbook for Market Domination", Inches(0.5), Inches(0.9), Inches(15), Inches(0.5), font_size=20, font_color=TEXT_COLOR_LIGHT)

# --- Main Content Area ---

# --- Left Column (Challenge & Opportunity) ---
add_textbox(slide, "The Challenge", Inches(0.5), Inches(1.8), Inches(3.5), Inches(0.5), font_size=22, bold=True, font_color=NEON_CYAN)
add_textbox(slide, "Design a scalable and profitable Tier-2 expansion plan that balances brand consistency with local adaptation and affordability.", Inches(0.5), Inches(2.3), Inches(3.5), Inches(1.5), font_size=14, font_color=TEXT_COLOR_LIGHT)

add_textbox(slide, "The Opportunity", Inches(0.5), Inches(3.8), Inches(3.5), Inches(0.5), font_size=22, bold=True, font_color=NEON_CYAN)
add_textbox(slide, "Tap into high-growth Tier-2 markets fueled by rising incomes and digital adoption, with lower AOV (₹200-220) but massive volume potential.", Inches(0.5), Inches(4.3), Inches(3.5), Inches(1.5), font_size=14, font_color=TEXT_COLOR_LIGHT)

# Projected Outcome Box
add_rounded_rect(slide, Inches(0.5), Inches(6.0), Inches(3.5), Inches(1.5), CARD_BG_COLOR, NEON_ORANGE)
add_textbox(slide, "Projected Outcome", Inches(0.6), Inches(6.1), Inches(3.3), Inches(0.4), font_size=18, bold=True, font_color=NEON_ORANGE)
add_textbox(slide, "Achieve kitchen-level profitability within 18-24 months, creating a proven playbook for national, asset-light expansion.", Inches(0.6), Inches(6.5), Inches(3.3), Inches(0.8), font_size=12, font_color=TEXT_COLOR_LIGHT)


# --- Right Grid (Strategic Pillars) ---
card_width = Inches(5.5)
card_height = Inches(2.2)

# Card 1: Launch Strategy
add_rounded_rect(slide, Inches(4.5), Inches(1.8), card_width, card_height, CARD_BG_COLOR, BORDER_COLOR)
add_textbox(slide, "Launch Strategy: Pilot First, Then Scale", Inches(4.7), Inches(2.0), Inches(5.1), Inches(0.5), font_size=18, bold=True, font_color=TEXT_COLOR_LIGHT)
add_textbox(slide, "De-risk expansion by launching in Indore first to test, learn, and create a replicable playbook before a wider rollout in Jaipur.", Inches(4.7), Inches(2.5), Inches(5.1), Inches(1.3), font_size=14, font_color=TEXT_COLOR_LIGHT)

# Card 2: Entry Model
add_rounded_rect(slide, Inches(10.2), Inches(1.8), card_width, card_height, CARD_BG_COLOR, BORDER_COLOR)
add_textbox(slide, "Entry Model: COCO to Franchise", Inches(10.4), Inches(2.0), Inches(5.1), Inches(0.5), font_size=18, bold=True, font_color=TEXT_COLOR_LIGHT)
add_textbox(slide, "Start with a Company-Owned (COCO) model for control, then transition to a Franchise model for rapid, asset-light growth.", Inches(10.4), Inches(2.5), Inches(5.1), Inches(1.3), font_size=14, font_color=TEXT_COLOR_LIGHT)

# Card 3: Adaptive Menu & Pricing
add_rounded_rect(slide, Inches(4.5), Inches(4.2), card_width, card_height, CARD_BG_COLOR, BORDER_COLOR)
add_textbox(slide, "Adaptive Menu & Pricing", Inches(4.7), Inches(4.3), Inches(3.0), Inches(0.4), font_size=18, bold=True, font_color=TEXT_COLOR_LIGHT)
# Pie Chart for Menu
chart_data = ChartData()
chart_data.categories = ['Core Menu (70%)', 'Local Menu (30%)']
chart_data.add_series('Menu Mix', (70, 30))
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(4.6), Inches(4.8), Inches(1.8), Inches(1.5), chart_data
).chart
chart.has_legend = False
chart.plots[0].has_data_labels = False
# Style Pie Chart Slices
point1 = chart.series[0].points[0]
point1.format.fill.solid()
point1.format.fill.fore_color.rgb = NEON_ORANGE
point2 = chart.series[0].points[1]
point2.format.fill.solid()
point2.format.fill.fore_color.rgb = NEON_CYAN
# Description text for pie chart
add_textbox(slide, "70% Core Menu: National 'hero' products for brand consistency & profitability.", Inches(6.5), Inches(4.8), Inches(3.3), Inches(0.8), font_size=11, font_color=TEXT_COLOR_LIGHT)
add_textbox(slide, "30% Local Menu: Hyper-local dishes to drive market penetration and trials.", Inches(6.5), Inches(5.6), Inches(3.3), Inches(0.8), font_size=11, font_color=TEXT_COLOR_LIGHT)


# Card 4: Go-To-Market
add_rounded_rect(slide, Inches(10.2), Inches(4.2), card_width, card_height, CARD_BG_COLOR, BORDER_COLOR)
add_textbox(slide, "Go-To-Market Strategy", Inches(10.4), Inches(4.3), Inches(5.1), Inches(0.4), font_size=18, bold=True, font_color=TEXT_COLOR_LIGHT)
g2m_text = "• Phase 1: Focus on aggregator data & retention.\n• Gamified Loyalty: Launch points-based rewards.\n• Local Collabs: LTOs with popular local chefs.\n• Phase 2: Scale with hyperlocal influencers."
add_textbox(slide, g2m_text, Inches(10.4), Inches(4.7), Inches(5.1), Inches(1.6), font_size=12, font_color=TEXT_COLOR_LIGHT)


# --- Footer (Projections Line Chart) ---
add_textbox(slide, "12-Month Projections (Indore Pilot)", Inches(4.5), Inches(6.6), Inches(11.2), Inches(0.5), font_size=20, bold=True, font_color=NEON_ORANGE, align='center')
# Line Chart Data
chart_data = ChartData()
chart_data.categories = ['M1', 'M2', 'M3', 'M4', 'M5', 'M6', 'M7', 'M8', 'M9', 'M10', 'M11', 'M12']
chart_data.add_series('Retention Rate (%)', (25, 26, 29, 27, 31, 33, 30, 34, 36, 34, 38, 41))

# Add and Style Line Chart
x, y, cx, cy = Inches(4.75), Inches(7.0), Inches(10.5), Inches(1.8)
graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)
line_chart = graphic_frame.chart

line_chart.has_legend = False
line_chart.chart_title.text_frame.text = 'Customer Retention Rate (%)'
line_chart.chart_title.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR_LIGHT

# --- FIX: Use the containing shape and plot_area to make chart background transparent ---
# `line_chart` is a Chart object (has no `.format` attribute like a shape does).
# The graphic_frame returned by add_chart is a shape and exposes `.fill`/`.format`.
# Clear both the graphic frame fill and the chart's plot_area fill (if available)
# to make the chart background transparent in a compatible way.
try:
    # Prefer clearing the chart area fill first
    line_chart.chart_area.format.fill.background()
except Exception:
    try:
        # Fallback to plot_area if chart_area isn't available
        line_chart.plot_area.format.fill.background()
    except Exception:
        # If neither is available, skip making the background transparent
        pass

# Style Axes
category_axis = line_chart.category_axis
category_axis.format.line.color.rgb = BORDER_COLOR
category_axis.tick_labels.font.color.rgb = TEXT_COLOR_LIGHT

value_axis = line_chart.value_axis
value_axis.format.line.color.rgb = BORDER_COLOR
value_axis.tick_labels.font.color.rgb = TEXT_COLOR_LIGHT
value_axis.has_major_gridlines = True
value_axis.major_gridlines.format.line.color.rgb = BORDER_COLOR

# Style Line Series
line_chart.series[0].format.line.color.rgb = NEON_ORANGE
line_chart.series[0].format.line.width = Pt(2.5)


# --- Save Presentation ---
file_path = 'QuickByte_Expansion_Strategy.pptx'
prs.save(file_path)

print(f"Presentation saved to {file_path}")

